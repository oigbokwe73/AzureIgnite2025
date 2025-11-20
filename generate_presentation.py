from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import os

# Create directories
os.makedirs('assets', exist_ok=True)
os.makedirs('presentation', exist_ok=True)

# Create a sample hero image with Pillow
img_path = os.path.join('assets', 'hero.png')
if not os.path.exists(img_path):
    W, H = 1280, 720
    img = Image.new('RGB', (W, H), color=(30, 144, 255))
    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype('DejaVuSans-Bold.ttf', 60)
    except Exception:
        font = ImageFont.load_default()
    text = 'Azure Ignite 2025'
    try:
        bbox = draw.textbbox((0, 0), text, font=font)
        w = bbox[2] - bbox[0]
        h = bbox[3] - bbox[1]
    except Exception:
        # fallback for older Pillow versions
        w, h = font.getsize(text)
    draw.text(((W-w)/2, (H-h)/2), text, fill=(255,255,255), font=font)
    img.save(img_path)

# Create presentation
prs = Presentation()

# Title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Azure Ignite 2025"
subtitle.text = "Sample slide deck generated with python-pptx"

# Image slide
img_slide_layout = prs.slide_layouts[5]  # title only
slide = prs.slides.add_slide(img_slide_layout)
slide.shapes.title.text = "Event Highlights"
left = Inches(1)
top = Inches(1.5)
width = Inches(8)
slide.shapes.add_picture(img_path, left, top, width=width)

# Bullet slide
bullet_slide_layout = prs.slide_layouts[1]  # title and content
slide = prs.slides.add_slide(bullet_slide_layout)
slide.shapes.title.text = "Key Takeaways"
body = slide.shapes.placeholders[1].text_frame
body.text = "- Cloud-native innovations"
p = body.add_paragraph()
p.text = "- Scalable architectures"
p.level = 1
p = body.add_paragraph()
p.text = "- Cost optimization strategies"

# Two-column content: quick features
two_col_layout = prs.slide_layouts[3]
slide = prs.slides.add_slide(two_col_layout)
slide.shapes.title.text = "Agenda"
left_tf = slide.shapes.placeholders[1].text_frame
left_tf.text = "1. Keynotes\n2. Workshops\n3. Demos"
right_tf = slide.shapes.placeholders[2].text_frame
right_tf.text = "4. Q&A\n5. Networking"

# Closing slide with contact and small logo
closing_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(closing_layout)
shapes = slide.shapes
if shapes.title:
    shapes.title.text = "Thank You"

left = Inches(0.5)
top = Inches(1.5)
width = Inches(3)
# add small copy of hero image as a logo
slide.shapes.add_picture(img_path, left, top, width=width)

# Save
out_path = os.path.join('presentation', 'sample_presentation.pptx')
prs.save(out_path)
print(f"Generated presentation: {out_path}")
